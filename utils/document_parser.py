# rag_chatbot/utils/document_parser.py

import os
from typing import List, Dict, Any

# Libraries for document parsing
# Ensure these are installed: pip install pypdf python-docx python-pptx pandas
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
    print("Warning: pypdf not installed. PDF parsing will not be available.")

try:
    from docx import Document
except ImportError:
    Document = None
    print("Warning: python-docx not installed. DOCX parsing will not be available.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("Warning: python-pptx not installed. PPTX parsing will not be available.")

try:
    import pandas as pd
except ImportError:
    pd = None
    print("Warning: pandas not installed. CSV parsing will not be available.")

def parse_pdf(file_path: str) -> str:
    """
    Parses a PDF file and extracts all text content.

    Args:
        file_path (str): The path to the PDF file.

    Returns:
        str: The concatenated text content from all pages of the PDF.

    Raises:
        ValueError: If pypdf is not installed or if the file cannot be read.
    """
    if PdfReader is None:
        raise ValueError("pypdf library is not installed. Cannot parse PDF files.")
    
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or "" # Use .extract_text() and handle None
        return text
    except Exception as e:
        raise ValueError(f"Error parsing PDF file {file_path}: {e}")

def parse_pptx(file_path: str) -> str:
    """
    Parses a PPTX file and extracts text from all slides.

    Args:
        file_path (str): The path to the PPTX file.

    Returns:
        str: The concatenated text content from all shapes in all slides.

    Raises:
        ValueError: If python-pptx is not installed or if the file cannot be read.
    """
    if Presentation is None:
        raise ValueError("python-pptx library is not installed. Cannot parse PPTX files.")

    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        raise ValueError(f"Error parsing PPTX file {file_path}: {e}")

def parse_csv(file_path: str) -> str:
    """
    Parses a CSV file and converts its content into a string representation.

    Args:
        file_path (str): The path to the CSV file.

    Returns:
        str: A string representation of the CSV data.

    Raises:
        ValueError: If pandas is not installed or if the file cannot be read.
    """
    if pd is None:
        raise ValueError("pandas library is not installed. Cannot parse CSV files.")
    
    try:
        df = pd.read_csv(file_path)
        # Convert DataFrame to a markdown-like table string or just its string representation
        return df.to_string(index=False)
    except Exception as e:
        raise ValueError(f"Error parsing CSV file {file_path}: {e}")

def parse_docx(file_path: str) -> str:
    """
    Parses a DOCX file and extracts all text content.

    Args:
        file_path (str): The path to the DOCX file.

    Returns:
        str: The concatenated text content from all paragraphs in the DOCX.

    Raises:
        ValueError: If python-docx is not installed or if the file cannot be read.
    """
    if Document is None:
        raise ValueError("python-docx library is not installed. Cannot parse DOCX files.")

    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        raise ValueError(f"Error parsing DOCX file {file_path}: {e}")

def parse_txt_md(file_path: str) -> str:
    """
    Reads a plain text or Markdown file and returns its content.

    Args:
        file_path (str): The path to the TXT/Markdown file.

    Returns:
        str: The full text content of the file.

    Raises:
        ValueError: If the file cannot be read.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        raise ValueError(f"Error reading TXT/Markdown file {file_path}: {e}")

def parse_document(file_path: str) -> str:
    """
    Parses a document based on its file extension and returns its text content.

    Args:
        file_path (str): The path to the document file.

    Returns:
        str: The extracted text content.

    Raises:
        ValueError: If the file format is not supported or parsing fails.
    """
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".pdf":
        return parse_pdf(file_path)
    elif file_extension == ".pptx":
        return parse_pptx(file_path)
    elif file_extension == ".csv":
        return parse_csv(file_path)
    elif file_extension == ".docx":
        return parse_docx(file_path)
    elif file_extension in [".txt", ".md"]:
        return parse_txt_md(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}. "
                         "Supported formats are .pdf, .pptx, .csv, .docx, .txt, .md.")

